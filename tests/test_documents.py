"""Document parsing.

Fixtures are generated rather than committed, so the tests exercise real files
produced by real writers instead of hand-rolled bytes that happen to parse.

The recurring theme is tables. Policy keeps its thresholds in them, the
grounding gate is built on figures being right, and a table flattened into prose
gives you numbers with no labels attached - which is the input that makes an
invented figure indistinguishable from a real one.
"""

from __future__ import annotations

from pathlib import Path

import pytest

from openknowledge.documents import (
    BlockKind,
    parse_bytes,
    parse_file,
    parse_text,
    skip_reason,
)
from openknowledge.documents.blocks import (
    Block,
    looks_like_header_row,
    normalise,
    table_row_text,
)


def kinds(doc) -> list[str]:
    return [b.kind.value for b in doc.blocks]


def texts(doc) -> list[str]:
    return [b.text for b in doc.blocks]


# -- blocks ----------------------------------------------------------------


def test_a_table_row_keeps_its_labels() -> None:
    """Without them the row is three unrelated numbers."""
    assert (
        table_row_text(["Grade", "Limit", "Notice"], ["Senior", "EUR 500", "60 days"])
        == "Grade: Senior | Limit: EUR 500 | Notice: 60 days"
    )


def test_unlabelled_columns_do_not_invent_a_header() -> None:
    assert table_row_text([], ["Senior", "EUR 500"]) == "Senior | EUR 500"


def test_empty_cells_are_dropped() -> None:
    assert table_row_text(["A", "B", "C"], ["x", "", "z"]) == "A: x | C: z"


@pytest.mark.parametrize(
    ("cells", "expected"),
    [
        (["Grade", "Limit", "Notice"], True),
        (["Senior", "EUR 500", "60 days"], False),  # has digits: data, not labels
        (["Grade"], False),  # one column proves nothing
        (["A very long sentence that is clearly prose and not a column label", "B"], False),
    ],
)
def test_header_row_detection_is_conservative(cells: list[str], expected: bool) -> None:
    """A wrong guess mislabels every row beneath it."""
    assert looks_like_header_row(cells) is expected


def test_contextual_text_carries_the_heading_trail() -> None:
    block = Block(BlockKind.PARAGRAPH, "Above EUR 500.", ("Expenses", "Thresholds"))
    assert block.contextual_text == "Expenses > Thresholds: Above EUR 500."


def test_a_heading_does_not_repeat_itself() -> None:
    block = Block(BlockKind.HEADING, "Thresholds", ("Expenses",))
    assert block.contextual_text == "Thresholds"


def test_normalise_collapses_extraction_damage() -> None:
    assert normalise("a  b\r\n\r\n\r\nc   d") == "a b\n\nc d"


def test_table_rows_are_atomic_and_paragraphs_are_not() -> None:
    assert BlockKind.TABLE_ROW.is_atomic
    assert not BlockKind.PARAGRAPH.is_atomic


# -- markdown / text -------------------------------------------------------

MARKDOWN = """# Expenses Policy

## Approval thresholds

Any expense above EUR 500 requires approval.

| Grade | Limit | Notice |
|---|---|---|
| Junior | EUR 200 | 5 days |
| Senior | EUR 1,000 | 2 days |

- Alcohol is never reimbursable
- Receipts are required above EUR 25
"""


def test_markdown_structure_is_recovered() -> None:
    doc = parse_text(MARKDOWN)
    assert doc.title == "Expenses Policy"
    assert kinds(doc) == [
        "heading",
        "heading",
        "paragraph",
        "table_row",
        "table_row",
        "list_item",
        "list_item",
    ]


def test_markdown_headings_nest() -> None:
    doc = parse_text(MARKDOWN)
    paragraph = next(b for b in doc.blocks if b.kind is BlockKind.PARAGRAPH)
    assert paragraph.heading_path == ("Expenses Policy", "Approval thresholds")


def test_markdown_table_rows_are_labelled() -> None:
    doc = parse_text(MARKDOWN)
    rows = [b.text for b in doc.blocks if b.kind is BlockKind.TABLE_ROW]
    assert rows[0] == "Grade: Junior | Limit: EUR 200 | Notice: 5 days"


def test_a_deeper_heading_does_not_escape_its_parent() -> None:
    doc = parse_text("# A\n\n## B\n\ntext\n\n# C\n\nmore")
    last = doc.blocks[-1]
    assert last.heading_path == ("C",), "a new h1 must reset the trail"


def test_setext_headings_are_recognised() -> None:
    doc = parse_text("Expenses Policy\n===\n\nSome text.")
    assert doc.blocks[0].kind is BlockKind.HEADING
    assert doc.title == "Expenses Policy"


def test_plain_text_without_structure_still_parses() -> None:
    doc = parse_text("Just one paragraph of prose.")
    assert kinds(doc) == ["paragraph"]


def test_empty_input_is_not_an_error() -> None:
    assert parse_text("   ").is_empty


# -- binary formats --------------------------------------------------------


@pytest.fixture
def corpus(tmp_path: Path) -> Path:
    """Real files, written by real writers."""
    from docx import Document as Docx
    from openpyxl import Workbook
    from pptx import Presentation
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    styles = getSampleStyleSheet()
    ruled = Table(
        [
            ["Grade", "Meal allowance", "Notice"],
            ["Junior", "EUR 35", "5 days"],
            ["Senior", "EUR 45", "2 days"],
        ]
    )
    ruled.setStyle(TableStyle([("GRID", (0, 0), (-1, -1), 0.5, colors.grey)]))
    SimpleDocTemplate(str(tmp_path / "policy.pdf"), pagesize=A4).build(
        [
            Paragraph("Expenses Policy", styles["Heading1"]),
            Paragraph("Approval thresholds", styles["Heading2"]),
            Paragraph(
                "Any single expense above EUR 500 requires prior written approval.",
                styles["Normal"],
            ),
            Spacer(1, 14),
            ruled,
        ]
    )

    doc = Docx()
    doc.add_heading("Remote Working", level=1)
    doc.add_heading("Eligibility", level=2)
    doc.add_paragraph("Employees with 6 months of service may work remotely.")
    table = doc.add_table(rows=2, cols=2)
    for r, row in enumerate([["Grade", "Days"], ["Senior", "3"]]):
        for c, value in enumerate(row):
            table.cell(r, c).text = value
    doc.add_heading("Equipment", level=2)
    doc.add_paragraph("Personal laptops are not permitted.")
    doc.save(tmp_path / "remote.docx")

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Caps"
    for row in [["Country", "Hotel cap"], ["Italy", 180], ["Japan", 240]]:
        sheet.append(row)
    workbook.save(tmp_path / "caps.xlsx")

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[1])
    slide.shapes.title.text = "How approvals work"
    slide.placeholders[1].text_frame.text = "Above EUR 500 needs your manager."
    deck.save(tmp_path / "approvals.pptx")

    return tmp_path


def test_pdf_recovers_headings_from_type_size(corpus: Path) -> None:
    doc = parse_file(corpus / "policy.pdf")
    assert doc.title == "Expenses Policy"
    assert "heading" in kinds(doc)


def test_pdf_keeps_tables_out_of_the_prose(corpus: Path) -> None:
    """The rows must appear once, as labelled rows - not again as a run-on line."""
    doc = parse_file(corpus / "policy.pdf")
    rows = [b.text for b in doc.blocks if b.kind is BlockKind.TABLE_ROW]
    assert any("Grade: Senior" in r and "EUR 45" in r for r in rows)

    prose = " ".join(b.text for b in doc.blocks if b.kind is BlockKind.PARAGRAPH)
    assert "EUR 45" not in prose, "table text was indexed twice"


def test_pdf_blocks_carry_a_page_locator(corpus: Path) -> None:
    """A citation an employee cannot check is decoration."""
    doc = parse_file(corpus / "policy.pdf")
    assert all(b.locator == "p. 1" for b in doc.blocks)
    assert doc.pages == 1


def test_pdf_prose_survives(corpus: Path) -> None:
    doc = parse_file(corpus / "policy.pdf")
    assert any("EUR 500" in b.text for b in doc.blocks if b.kind is BlockKind.PARAGRAPH)


def test_docx_keeps_document_order(corpus: Path) -> None:
    """python-docx exposes paragraphs and tables as two flat lists; relying on
    them puts every table after all the prose, detached from its heading."""
    doc = parse_file(corpus / "remote.docx")
    order = kinds(doc)
    table_at = order.index("table_row")
    equipment_at = next(i for i, b in enumerate(doc.blocks) if b.text == "Equipment")
    assert table_at < equipment_at, "the table must precede the heading that follows it"


def test_docx_headings_and_tables(corpus: Path) -> None:
    doc = parse_file(corpus / "remote.docx")
    assert doc.title == "Remote Working"
    row = next(b for b in doc.blocks if b.kind is BlockKind.TABLE_ROW)
    assert row.text == "Grade: Senior | Days: 3"
    assert row.heading_path == ("Remote Working", "Eligibility")


def test_xlsx_rows_are_labelled_and_cell_addressed(corpus: Path) -> None:
    doc = parse_file(corpus / "caps.xlsx")
    rows = [b for b in doc.blocks if b.kind is BlockKind.TABLE_ROW]
    assert rows[0].text == "Country: Italy | Hotel cap: 180"
    assert rows[0].locator == "Caps!A2"


def test_pptx_uses_slide_titles_and_numbers(corpus: Path) -> None:
    doc = parse_file(corpus / "approvals.pptx")
    assert doc.title == "How approvals work"
    assert all(b.locator and b.locator.startswith("slide") for b in doc.blocks)


# -- failure handling ------------------------------------------------------


def test_a_corrupt_file_reports_rather_than_raises(tmp_path: Path) -> None:
    """One bad file must not take down a whole re-index."""
    bad = tmp_path / "broken.pdf"
    bad.write_bytes(b"%PDF-1.4 this is not really a pdf")
    doc = parse_file(bad)
    assert doc.is_empty
    assert doc.warnings


def test_a_pdf_with_no_text_layer_is_named_as_a_scan(tmp_path: Path) -> None:
    """Silently indexing nothing is how a corpus grows a hole nobody notices."""
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas

    path = tmp_path / "scan.pdf"
    c = canvas.Canvas(str(path), pagesize=A4)
    c.rect(100, 100, 200, 200, fill=0)  # a drawing, no text
    c.save()

    doc = parse_file(path)
    assert any("scan" in w for w in doc.warnings)
    assert any("OCR" in w for w in doc.warnings)


@pytest.mark.parametrize(
    ("name", "fragment"),
    [
        ("old.doc", "re-save it as .docx"),
        ("old.xls", "re-save it as .xlsx"),
        ("image.png", "no parser for .png"),
    ],
)
def test_unreadable_formats_say_what_to_do(name: str, fragment: str) -> None:
    reason = skip_reason(name)
    assert reason and fragment in reason


def test_a_supported_format_has_no_skip_reason() -> None:
    assert skip_reason("policy.pdf") is None


def test_parse_bytes_dispatches_on_suffix() -> None:
    doc = parse_bytes(b"# Title\n\nBody.", suffix=".md")
    assert doc.title == "Title"
    assert parse_bytes(b"x", suffix=".bin").warnings
