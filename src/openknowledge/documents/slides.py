"""Slide decks.

Internal process is often documented in a deck - an onboarding walkthrough, a
"how approvals work" pack - so skipping the format leaves a real hole. Each
slide's title becomes the heading, everything else hangs under it, and the
locator is the slide number, which is how people refer to decks anyway.

Speaker notes are included: on a process deck they frequently hold the actual
detail the slide only gestures at.
"""

from __future__ import annotations

import contextlib
import logging

from .blocks import (
    Block,
    BlockKind,
    ParsedDocument,
    looks_like_header_row,
    normalise,
    table_row_text,
)

log = logging.getLogger(__name__)


def parse_pptx(data: bytes, *, title: str | None = None) -> ParsedDocument:
    """Parse a .pptx into per-slide blocks."""
    try:
        from pptx import Presentation
    except ImportError:  # pragma: no cover - dependency is declared
        return ParsedDocument(warnings=("python-pptx is not installed",))

    import io

    try:
        deck = Presentation(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        log.warning("could not read deck: %s", exc)
        return ParsedDocument(warnings=(f"could not read this presentation: {exc}",))

    blocks: list[Block] = []
    doc_title = title

    slide_count = 0
    for number, slide in enumerate(deck.slides, start=1):
        slide_count = number
        locator = f"slide {number}"
        slide_title = ""
        try:
            if slide.shapes.title is not None:
                slide_title = normalise(slide.shapes.title.text)
        except Exception:  # noqa: BLE001 - layouts without a title placeholder
            slide_title = ""

        heading = slide_title or f"Slide {number}"
        if doc_title is None and number == 1 and slide_title:
            doc_title = slide_title
        blocks.append(Block(kind=BlockKind.HEADING, text=heading, locator=locator, level=1))
        path = (heading,)

        for shape in slide.shapes:
            if shape.has_table:
                blocks.extend(_table_blocks(shape.table, path, locator))
                continue
            if not shape.has_text_frame:
                continue
            if slide_title and shape == slide.shapes.title:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = normalise("".join(run.text for run in paragraph.runs))
                if text:
                    blocks.append(
                        Block(
                            kind=BlockKind.LIST_ITEM if paragraph.level else BlockKind.PARAGRAPH,
                            text=text,
                            heading_path=path,
                            locator=locator,
                        )
                    )

        # Decks without a notes master raise rather than reporting no notes.
        with contextlib.suppress(Exception):
            if slide.has_notes_slide:
                notes = normalise(slide.notes_slide.notes_text_frame.text)
                if notes:
                    blocks.append(
                        Block(
                            kind=BlockKind.CAPTION,
                            text=notes,
                            heading_path=path,
                            locator=f"{locator} notes",
                        )
                    )

    return ParsedDocument(blocks=tuple(blocks), title=doc_title, pages=slide_count)


def _table_blocks(table: object, path: tuple[str, ...], locator: str) -> list[Block]:
    blocks: list[Block] = []
    headers: list[str] = []
    for row in table.rows:  # type: ignore[attr-defined]
        cells = [normalise(cell.text) for cell in row.cells]
        if not headers and looks_like_header_row(cells):
            headers = cells
            continue
        text = table_row_text(headers, cells)
        if text:
            blocks.append(
                Block(kind=BlockKind.TABLE_ROW, text=text, heading_path=path, locator=locator)
            )
    return blocks
